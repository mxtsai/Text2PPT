# Designed by LiteHouse.Taipei
# Use this script to speed up PPT development for worship purposes
# Modify 'template.pptx' to change the font, text size (Open in PPT > View > Slide Master > and modify)

import os
from pptx import Presentation

def extract_lyrics(lyrics_text_file, lines_per_slide = 1):
    lyrics = []
    try:
        with open(lyrics_text_file) as f:
            line = f.readlines()    # read raw lines into a list
            
            cur_line_count = 1
            cur_line_txt = ""   # to append any text
            for l in line:
                line_text = l.strip()

                if line_text:   # if have string content
                    if cur_line_count >= lines_per_slide:   # if have enough lines
                        if cur_line_txt.strip():    # there exists prev line
                            lyrics.append(f"{cur_line_txt}\n{line_text}")   # add to lyrics
                        else:
                            lyrics.append(line_text)    # add only this line to lyrics
                        cur_line_count = 1  # reset line count
                        cur_line_txt = ""   # reset line text
                    else:   # if not enough lines
                        cur_line_txt = line_text    # add to line text buffer
                        cur_line_count += 1     # increment line count
                else:   # empty line
                    if cur_line_txt.strip():    # if exists string in buffer
                        lyrics.append(f"{cur_line_txt}")    # add it to lyrics
                    cur_line_count = 1  # reset line count
                    cur_line_txt = ""   # reset line text
            if cur_line_txt.strip():    # if exists string in buffer
                lyrics.append(f"{cur_line_txt}")    # add it to lyrics
    except:
        print(f"Error: cannot read file {lyrics_text_file}")

    return lyrics


def generate_lyrics_ppt(lyrics_list, save_path):
    """takes a list of lyrics (inclues formatted string) and generates a pptx file"""

    lyrics_slides = Presentation('template.pptx')
    layout = lyrics_slides.slide_master.slide_layouts[0]

    # add slides to PPTX
    for line in lyrics_list:
        slide = lyrics_slides.slides.add_slide(layout)
        slide.shapes.title.text = line

    # save 
    os.makedirs(os.path.dirname(save_path), exist_ok=True)
    lyrics_slides.save(save_path)


if __name__ == "__main__":

    input_message = """
This program generates a pptx file from a text file of lyrics.
Drag the txt file with the lyrics to this program, and it will generate a pptx file with the same name. (press Ctrl-C to cancel).
"""

    try:
        
        song_file = input(input_message).strip().lstrip('\"').rstrip('\"').replace("\\", "")
        song_file = os.path.abspath(song_file)

        if song_file:
            lines_per_slide = int(input("\nPlease specify the number of lines to print on each slide:\n").strip())

            file_dir, file_name = os.path.split(song_file)

            # extract lyrics
            lyrics_list = extract_lyrics(song_file, lines_per_slide=lines_per_slide)

            # generate PPT
            filename_out = (file_name.split('.')[0])
            save_path = f"{file_dir}/{filename_out}.pptx"
            generate_lyrics_ppt(lyrics_list, save_path)
            print("PPT successfully generated. Check the file at:", save_path)

    except Exception as e:
        print(e)
    input("Press any key to exit\n")


